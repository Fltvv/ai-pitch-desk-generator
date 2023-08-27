import os

import collections
import collections.abc

from pptx import Presentation
from pptx.util import Inches

from pptx.dml.color import RGBColor

from connection_to_sql import SqlClass

from random import choice


sql = SqlClass()

picture_types = ['description', 'market', 'name', 'product']


def get_number_of_slides(presentation):
    number_of_slides = len(presentation.slides)
    return number_of_slides


def check_is_text_shape(shape):
    is_there_text = False
    if shape.has_text_frame:
        if len(shape.text_frame.text) > 0:
            is_there_text = True
    return is_there_text


def get_text_shapes_indexes(slide):
    list_of_shapes = list(slide.shapes)
    text_shapes_indexes = []
    for i, shape in enumerate(list_of_shapes):
        if check_is_text_shape(shape):
            text_shapes_indexes.append(i)
    return text_shapes_indexes


def get_text_shapes(slide):
    list_of_shapes = list(slide.shapes)
    text_shapes_indexes = get_text_shapes_indexes(slide)
    text_shapes = [list_of_shapes[i] for i in text_shapes_indexes]
    return text_shapes


def get_text(text_shape):
    text = text_shape.text_frame.text
    return text


def set_text_color_and_font(presentation,
                            slide_number,
                            shape_number,
                            background_style,
                            font_name):
    number_of_paragraphs = len(presentation.slides[slide_number].shapes[shape_number].text_frame.paragraphs)
    for number_of_paragraph in range(number_of_paragraphs):
        presentation.slides[slide_number].shapes[shape_number].text_frame.paragraphs[number_of_paragraph].font.name = font_name
        if background_style == 'white':
            presentation.slides[slide_number].shapes[shape_number].text_frame.paragraphs[number_of_paragraph].font.color.rgb = RGBColor(0, 0, 0)

            num_iterations = len(presentation.slides[slide_number].shapes[shape_number].text_frame.paragraphs[number_of_paragraph].runs)
            for k in range(num_iterations):
                presentation.slides[slide_number].shapes[shape_number].text_frame.paragraphs[number_of_paragraph].runs[k].font.color.rgb = RGBColor(0, 0, 0)

        elif background_style == 'black':
            presentation.slides[slide_number].shapes[shape_number].text_frame.paragraphs[number_of_paragraph].font.color.rgb = RGBColor(255, 255, 255)

            num_iterations = len(presentation.slides[slide_number].shapes[shape_number].text_frame.paragraphs[number_of_paragraph].runs)
            for k in range(num_iterations):
                presentation.slides[slide_number].shapes[shape_number].text_frame.paragraphs[number_of_paragraph].runs[k].font.color.rgb = RGBColor(255, 255, 255)

    return presentation


def set_text_with_required_properties(presentation,
                                      slide_number,
                                      shape_number,
                                      old_text_shape,
                                      new_text):
    alignment = old_text_shape.text_frame.paragraphs[0].alignment
    font_size = old_text_shape.text_frame.paragraphs[0].runs[0].font.size
    font_bold = old_text_shape.text_frame.paragraphs[0].runs[0].font.bold
    font_italic = old_text_shape.text_frame.paragraphs[0].runs[0].font.italic

    presentation.slides[slide_number].shapes[shape_number].text_frame.text = new_text

    number_of_paragraphs = len(presentation.slides[slide_number].shapes[shape_number].text_frame.paragraphs)
    for number_of_paragraph in range(number_of_paragraphs):
        presentation.slides[slide_number].shapes[shape_number].text_frame.paragraphs[number_of_paragraph].alignment = alignment
        presentation.slides[slide_number].shapes[shape_number].text_frame.paragraphs[number_of_paragraph].font.size = font_size
        presentation.slides[slide_number].shapes[shape_number].text_frame.paragraphs[number_of_paragraph].font.bold = font_bold
        presentation.slides[slide_number].shapes[shape_number].text_frame.paragraphs[number_of_paragraph].font.italic = font_italic

    return presentation


def change_text_on_slide(presentation,
                         slide_number,
                         id_startup,
                         language='ru',
                         font_name='Roboto',
                         background_style='white'):
    slide = presentation.slides[slide_number]
    list_of_shapes = list(slide.shapes)
    text_shapes_indexes = get_text_shapes_indexes(slide)

    for shape_number in text_shapes_indexes:
        old_text_shape = list_of_shapes[shape_number]
        text = get_text(old_text_shape)

        db_column_names = sql.get_base_column_names()

        if 'url_ppt' in db_column_names:
            db_column_names.remove('url_ppt')
        else:
            pass

        for column_name in db_column_names:
            if text == column_name:

                if language == 'ru':
                    new_text = sql.get_value_from_base(column_name, id_startup=id_startup)
                elif language == 'en':
                    new_text = sql.get_value_from_mybase(column_name, id_startup=id_startup)
                else:
                    new_text = sql.get_value_from_base(column_name, id_startup=id_startup)

                if len(new_text.split()) > 0:
                    pass
                else:
                    new_text = f'NO TEXT IN DB FOR COLUMN {column_name}'

                presentation = set_text_with_required_properties(presentation=presentation,
                                                                 new_text=new_text,
                                                                 old_text_shape=old_text_shape,
                                                                 slide_number=slide_number,
                                                                 shape_number=shape_number,
                                                                 )
            else:
                pass

        presentation = set_text_color_and_font(presentation=presentation,
                                               slide_number=slide_number,
                                               shape_number=shape_number,
                                               background_style=background_style,
                                               font_name=font_name
                                               )

    return presentation


def set_picture_with_required_properties(presentation,
                                         old_text_shape,
                                         slide_number,
                                         shape_number,
                                         path_to_image):
    top = old_text_shape.top
    left = old_text_shape.left
    width = old_text_shape.width.inches
    height = old_text_shape.height.inches

    remove_shape(presentation.slides[slide_number].shapes[shape_number])

    try:
        presentation.slides[slide_number].shapes.add_picture(path_to_image,
                                                             left,
                                                             top,
                                                             width=Inches(width),
                                                             height=Inches(height)
                                                             )
    except:
        pass

    return None


def create_picture_on_slide(presentation,
                            slide_number,
                            id_startup):
    slide = presentation.slides[slide_number]
    list_of_shapes = list(slide.shapes)
    text_shapes_indexes = get_text_shapes_indexes(slide)
    for i in text_shapes_indexes:
        old_text_shape = list_of_shapes[i]
        text = get_text(old_text_shape)

        elements = text.split('_')

        if (elements[0] == 'pic') and (elements[-1] in picture_types):
            picture_type = elements[-1]

            path_to_image = f'./images/{id_startup}_{picture_type}.jpg'

            set_picture_with_required_properties(presentation=presentation,
                                                 old_text_shape=old_text_shape,
                                                 slide_number=slide_number,
                                                 shape_number=i,
                                                 path_to_image=path_to_image
                                                 )
    return presentation


def create_background(presentation,
                      slide_number,
                      slide_width,
                      slide_height,
                      path_to_background_style):
    background_images = os.listdir(path_to_background_style)

    chosen_background_image = choice(background_images)

    path_to_background = path_to_background_style + f'/{chosen_background_image}'

    slide = presentation.slides[slide_number]

    left = Inches(0)
    top = Inches(0)

    slide_width = Inches(slide_width)
    slide_height = Inches(slide_height)

    pic = slide.shapes.add_picture(path_to_background,
                                   left,
                                   top,
                                   width=slide_width,
                                   height=slide_height
                                   )

    slide.shapes._spTree.remove(pic._element)
    slide.shapes._spTree.insert(2, pic._element)

    return presentation


def remove_shape(shape):
    shape = shape._element
    shape.getparent().remove(shape)
    return None


def create_presentation(template_presentation,
                        id_startup,
                        background_style='black',
                        language='ru',
                        save_as='template_changed.pptx'):
    slide_width = template_presentation.slide_width.inches
    slide_height = template_presentation.slide_height.inches

    if background_style == 'black':
        path_to_background_style = './backgrounds/black'
    elif background_style == 'white':
        path_to_background_style = './backgrounds/white'
    else:
        '''In other cases the background type is black'''
        path_to_background_style = './backgrounds/black'

    background_types = os.listdir(path_to_background_style)
    if '.DS_Store' in background_types:
        background_types.remove('.DS_Store')

    chosen_background_type = choice(background_types)

    path_to_background_style = path_to_background_style + f'/{chosen_background_type}'

    number_of_slides = get_number_of_slides(presentation=template_presentation)
    for slide_number in range(number_of_slides):

        template_presentation = change_text_on_slide(presentation=template_presentation,
                                                     slide_number=slide_number,
                                                     id_startup=id_startup,
                                                     language=language,
                                                     font_name='Roboto',
                                                     background_style=background_style
                                                     )
        # print(f'Text was changed on slide №{slide_number + 1}')

        template_presentation = create_background(presentation=template_presentation,
                                                  slide_number=slide_number,
                                                  slide_width=slide_width,
                                                  slide_height=slide_height,
                                                  path_to_background_style=path_to_background_style
                                                  )
        # print(f'Background was changed on slide №{slide_number + 1}')

        template_presentation = create_picture_on_slide(presentation=template_presentation,
                                                        slide_number=slide_number,
                                                        id_startup=id_startup)

        # print(f'Picture was added on slide №{slide_number + 1}')

    print('PITCH-DESK IS CREATED')

    template_presentation.save(save_as)

    return template_presentation
